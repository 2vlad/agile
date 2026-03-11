import csv
import hashlib
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path

import fitz  # PyMuPDF
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)


class SkippedDocument(Exception):
    """Raised when a document cannot be meaningfully parsed (e.g. scan-only PDF)."""


@dataclass
class ParsedDocument:
    filename: str
    format: str
    title: str
    text: str
    metadata: dict = field(default_factory=dict)
    content_hash: str = ""

    def __post_init__(self) -> None:
        if not self.content_hash:
            self.content_hash = hashlib.sha256(self.text.encode()).hexdigest()


SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".epub", ".md", ".csv", ".docx", ".pptx", ".xlsx", ".html", ".htm"}


def parse_file(filepath: Path) -> ParsedDocument:
    """Dispatch to correct parser based on file extension."""
    ext = filepath.suffix.lower()
    if ext in (".txt", ".md"):
        return parse_txt(filepath)
    if ext == ".pdf":
        return parse_pdf(filepath)
    if ext == ".epub":
        return parse_epub(filepath)
    if ext == ".docx":
        return parse_docx(filepath)
    if ext == ".pptx":
        return parse_pptx(filepath)
    if ext == ".xlsx":
        return parse_xlsx(filepath)
    if ext == ".csv":
        return parse_csv(filepath)
    if ext in (".html", ".htm"):
        return parse_html(filepath)
    raise ValueError(f"Unsupported file format: {ext}")


def parse_txt(filepath: Path) -> ParsedDocument:
    """Parse a plain text file. Try UTF-8 then cp1251 (common for Russian texts)."""
    raw = filepath.read_bytes()
    content_hash = hashlib.sha256(raw).hexdigest()

    text = None
    for encoding in ("utf-8", "cp1251", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue

    if text is None:
        raise ValueError(f"Cannot decode {filepath.name}")

    return ParsedDocument(
        filename=filepath.name,
        format="txt",
        title=filepath.stem.replace("_", " ").title(),
        text=text.strip(),
        content_hash=content_hash,
    )


def parse_pdf(filepath: Path) -> ParsedDocument:
    """Parse a text-based PDF page by page using PyMuPDF.

    Raises SkippedDocument if the PDF appears to be a scan (no extractable text).
    """
    content_hash = hashlib.sha256(filepath.read_bytes()).hexdigest()

    with fitz.open(str(filepath)) as doc:
        total_pages = doc.page_count
        meta = doc.metadata or {}

        pages_text: list[str] = []
        pages_with_text = 0

        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                pages_with_text += 1
                pages_text.append(f"[Страница {page_num}]\n{page_text}")

    if pages_with_text == 0:
        raise SkippedDocument(f"{filepath.name}: no extractable text (likely a scan)")

    # Extract title from PDF metadata, fallback to filename
    title = (meta.get("title") or "").strip() or filepath.stem.replace("_", " ").title()

    text = "\n\n".join(pages_text)
    return ParsedDocument(
        filename=filepath.name,
        format="pdf",
        title=title,
        text=text,
        metadata={"total_pages": total_pages},
        content_hash=content_hash,
    )


def parse_epub(filepath: Path) -> ParsedDocument:
    """Parse an EPUB file, extracting text from HTML spine items chapter by chapter."""
    content_hash = hashlib.sha256(filepath.read_bytes()).hexdigest()

    book = epub.read_epub(str(filepath), options={"ignore_ncx": True})

    # Extract title from EPUB metadata
    title_meta = book.get_metadata("DC", "title")
    title = title_meta[0][0] if title_meta else filepath.stem.replace("_", " ").title()

    chapters: list[str] = []
    chapter_titles: list[str] = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")

        # Extract chapter title from headings
        heading = soup.find(["h1", "h2", "h3"])
        chapter_title = heading.get_text(strip=True) if heading else ""

        # Remove scripts and styles
        for tag in soup(["script", "style", "head"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        text = "\n".join(line for line in text.splitlines() if line.strip())

        if len(text) > 100:  # Skip nearly-empty chapters
            if chapter_title:
                chapters.append(f"[{chapter_title}]\n{text}")
                chapter_titles.append(chapter_title)
            else:
                chapters.append(text)

    if not chapters:
        raise SkippedDocument(f"{filepath.name}: no text content found in EPUB")

    full_text = "\n\n".join(chapters)
    return ParsedDocument(
        filename=filepath.name,
        format="epub",
        title=title,
        text=full_text,
        metadata={"chapters": chapter_titles, "chapter_count": len(chapters)},
        content_hash=content_hash,
    )


def parse_docx(filepath: Path) -> ParsedDocument:
    """Parse a .docx file extracting text from paragraphs."""
    from docx import Document

    content_hash = hashlib.sha256(filepath.read_bytes()).hexdigest()
    doc = Document(str(filepath))

    title = (doc.core_properties.title or "").strip() or filepath.stem.replace("_", " ").title()

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        raise SkippedDocument(f"{filepath.name}: no text content in DOCX")

    return ParsedDocument(
        filename=filepath.name,
        format="docx",
        title=title,
        text="\n\n".join(paragraphs),
        content_hash=content_hash,
    )


def parse_pptx(filepath: Path) -> ParsedDocument:
    """Parse a .pptx file extracting text from slides and notes."""
    from pptx import Presentation

    content_hash = hashlib.sha256(filepath.read_bytes()).hexdigest()
    prs = Presentation(str(filepath))

    title = filepath.stem.replace("_", " ").title()
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Заметки] {notes}")
        if parts:
            slides_text.append(f"[Слайд {i}]\n" + "\n".join(parts))

    # Try to get title from first slide
    if prs.slides and prs.slides[0].shapes.title:
        first_title = prs.slides[0].shapes.title.text.strip()
        if first_title:
            title = first_title

    if not slides_text:
        raise SkippedDocument(f"{filepath.name}: no text content in PPTX")

    return ParsedDocument(
        filename=filepath.name,
        format="pptx",
        title=title,
        text="\n\n".join(slides_text),
        metadata={"slide_count": len(prs.slides)},
        content_hash=content_hash,
    )


def parse_xlsx(filepath: Path) -> ParsedDocument:
    """Parse an .xlsx file — each row as a line, sheets separated."""
    from openpyxl import load_workbook

    content_hash = hashlib.sha256(filepath.read_bytes()).hexdigest()
    wb = load_workbook(str(filepath), read_only=True, data_only=True)

    sheets_text: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets_text.append(f"[{sheet.title}]\n" + "\n".join(rows))

    wb.close()

    if not sheets_text:
        raise SkippedDocument(f"{filepath.name}: no data in XLSX")

    return ParsedDocument(
        filename=filepath.name,
        format="xlsx",
        title=filepath.stem.replace("_", " ").title(),
        text="\n\n".join(sheets_text),
        content_hash=content_hash,
    )


def parse_csv(filepath: Path) -> ParsedDocument:
    """Parse a .csv file — rows as pipe-separated lines."""
    raw = filepath.read_bytes()
    content_hash = hashlib.sha256(raw).hexdigest()

    text = None
    for encoding in ("utf-8", "cp1251", "latin-1"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        raise ValueError(f"Cannot decode {filepath.name}")

    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(cell for cell in row if cell.strip()) for row in reader]
    rows = [r for r in rows if r.strip()]

    if not rows:
        raise SkippedDocument(f"{filepath.name}: no data in CSV")

    return ParsedDocument(
        filename=filepath.name,
        format="csv",
        title=filepath.stem.replace("_", " ").title(),
        text="\n".join(rows),
        content_hash=content_hash,
    )


def parse_html(filepath: Path) -> ParsedDocument:
    """Parse an HTML file extracting text content."""
    raw = filepath.read_bytes()
    content_hash = hashlib.sha256(raw).hexdigest()

    soup = BeautifulSoup(raw, "html.parser")

    title_tag = soup.find("title")
    title = title_tag.get_text(strip=True) if title_tag else filepath.stem.replace("_", " ").title()

    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    text = "\n".join(line for line in text.splitlines() if line.strip())

    if len(text) < 50:
        raise SkippedDocument(f"{filepath.name}: too little text in HTML")

    return ParsedDocument(
        filename=filepath.name,
        format="html",
        title=title,
        text=text,
        content_hash=content_hash,
    )
